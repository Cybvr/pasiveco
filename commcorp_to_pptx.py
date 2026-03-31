"""
CommCorp HTML Deck → PPTX Converter
Faithfully recreates the 7-slide deck using python-pptx.
Output: commcorp_deck.pptx (same directory as this script)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

# ── Palette ──────────────────────────────────────────────────────────────────
RED   = RGBColor(0xCC, 0x15, 0x00)
DARK  = RGBColor(0x0E, 0x0E, 0x0E)
MID   = RGBColor(0x1A, 0x1A, 0x1A)
BONE  = RGBColor(0xF5, 0xF0, 0xE8)
SAND  = RGBColor(0xE8, 0xDF, 0xD0)
DIM   = RGBColor(0x88, 0x88, 0x88)
INK   = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY6 = RGBColor(0x66, 0x66, 0x66)
GREY9 = RGBColor(0x99, 0x99, 0x99)

# ── Slide dimensions: 16:9 ───────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank layout

# ── Helper: add a filled rectangle (no text) ────────────────────────────────
def rect(slide, x, y, w, h, fill_color, border=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if border:
        shape.line.color.rgb = SAND
        shape.line.width = Pt(0.5)
    return shape

# ── Helper: add a text box ───────────────────────────────────────────────────
def tb(slide, text, x, y, w, h,
       font_name="DM Sans", font_size=12, bold=False, italic=False,
       color=WHITE, align=PP_ALIGN.LEFT, wrap=True, auto_size=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name      = font_name
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    return txb

# ── Helper: multi-paragraph text box ────────────────────────────────────────
def mtb(slide, lines, x, y, w, h,
        font_name="DM Sans", font_size=10, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT, spacing_after=Pt(4)):
    """lines = list of (text, override_dict) or just strings"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if isinstance(line, str):
            text, overrides = line, {}
        else:
            text, overrides = line[0], line[1]

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = overrides.get("align", align)
        p.space_after = overrides.get("spacing_after", spacing_after)

        run = p.add_run()
        run.text = text
        run.font.name      = overrides.get("font_name",  font_name)
        run.font.size      = Pt(overrides.get("font_size",  font_size))
        run.font.bold      = overrides.get("bold",  bold)
        run.font.italic    = overrides.get("italic", italic)
        run.font.color.rgb = overrides.get("color",  color)
    return txb

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1: TITLE
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)

half = W / 2

# Left red panel
rect(s1, 0, 0, half, H, RED)

# Logo
tb(s1, "COMM\nCORP", Inches(0.4), Inches(0.4), half - Inches(0.5), Inches(3),
   font_name="Arial Black", font_size=62, bold=True, color=WHITE)

# Eyebrow
tb(s1, "Creative Business Partners — Lagos & Beyond",
   Inches(0.4), H - Inches(1.1), half - Inches(0.5), Inches(0.6),
   font_size=7.5, color=RGBColor(0xFF,0xFF,0xFF), bold=False)

# Right dark panel
rect(s1, half, 0, half, H, DARK)

# Tagline
tb(s1,
   "Nigeria's enterprise teams are building world-class brands.\nThey need world-class creative execution to match.",
   half + Inches(0.3), Inches(0.6), half - Inches(0.5), Inches(1.5),
   font_name="Georgia", font_size=11, italic=True, color=WHITE)

# Stats
for i, (num, lbl) in enumerate([
    ("₦4.7T", "Nigeria's advertising & marketing services industry, growing 12% year-on-year"),
    ("68%",   "of Nigerian CMOs cite talent mismatch as their top creative execution challenge"),
]):
    yy = Inches(2.5) + i * Inches(2.2)
    tb(s1, num, half + Inches(0.3), yy, half - Inches(0.5), Inches(0.9),
       font_name="Arial Black", font_size=36, bold=True, color=RED)
    tb(s1, lbl, half + Inches(0.3), yy + Inches(0.85), half - Inches(0.5), Inches(0.8),
       font_size=8, color=DIM)

# Page number
tb(s1, "01 / 07", W - Inches(1), H - Inches(0.3), Inches(0.8), Inches(0.25),
   font_size=7, color=RGBColor(0x55,0x55,0x55), align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2: THE PROBLEM
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
rect(s2, 0, 0, W, H, BONE)

# Header bar
hbar_h = Inches(0.95)
rect(s2, 0, 0, W, hbar_h, INK)
tb(s2, "The Problem Is Real", Inches(0.4), Inches(0.15), Inches(6), Inches(0.7),
   font_name="Arial Black", font_size=26, bold=True, color=WHITE)
tb(s2, "and it costs more than money", Inches(6.8), Inches(0.3), Inches(4), Inches(0.5),
   font_size=8, color=DIM)

# Three columns
col_w = W / 3
problems = [
    ("1", "The Ghost Vendor",
     "A Nigerian brand manager spends 3 weeks briefing a freelancer. They go dark after the advance drops. The campaign misses its window. It happens constantly — and there's no structure to prevent it.",
     "Pattern observed across 40+ enterprise engagements in Lagos, Abuja & PH"),
    ("2", "The Import Default",
     "Brands assume global agencies understand local context. They don't. Ads that land flat in Yaba, copy that misses pidgin tonality, visuals that say nothing to a Kano audience. Expensive mistakes with polished invoices.",
     "Deloitte Africa Consumer Survey, 2023"),
    ("3", "The Accountability Gap",
     "Once creative is outsourced, the internal team loses control. Revisions become arguments. Timelines become negotiations. And the marketing director takes the heat for work they couldn't actually direct.",
     "McKinsey Africa Creative Economy Report"),
]

for i, (num, title, body, src) in enumerate(problems):
    cx = i * col_w
    cy = hbar_h
    ch = H - hbar_h

    # Column background (alternating subtle tint)
    bg = RGBColor(0xF0, 0xEB, 0xE0) if i % 2 == 1 else BONE
    rect(s2, cx, cy, col_w, ch, bg)

    # Big ghost number
    tb(s2, num, cx + col_w - Inches(1), cy + Inches(0.1), Inches(0.9), Inches(1.2),
       font_name="Arial Black", font_size=54, bold=True, color=SAND, align=PP_ALIGN.RIGHT)

    # Title in red
    tb(s2, title.upper(), cx + Inches(0.25), cy + Inches(1.1), col_w - Inches(0.4), Inches(0.5),
       font_size=8.5, bold=True, color=RED)

    # Body
    tb(s2, body, cx + Inches(0.25), cy + Inches(1.65), col_w - Inches(0.4), Inches(3.5),
       font_size=8, color=RGBColor(0x44,0x44,0x44))

    # Source
    tb(s2, src, cx + Inches(0.25), H - Inches(0.55), col_w - Inches(0.4), Inches(0.4),
       font_size=7, italic=True, color=RGBColor(0xAA,0xAA,0xAA))

    # Divider line between columns
    if i < 2:
        divider = s2.shapes.add_shape(1, cx + col_w - Pt(0.25), cy, Pt(0.5), ch)
        divider.fill.solid()
        divider.fill.fore_color.rgb = SAND
        divider.line.fill.background()

tb(s2, "02 / 07", W - Inches(1), H - Inches(0.3), Inches(0.8), Inches(0.25),
   font_size=7, color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3: THE DIFFERENCE
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
rect(s3, 0, 0, W, H, DARK)

left_w = W * 45 // 100

# Red left panel
rect(s3, 0, 0, left_w, H, RED)

# Kicker
tb(s3, "THE COMMCORP DIFFERENCE", Inches(0.4), Inches(0.4), left_w - Inches(0.5), Inches(0.4),
   font_size=7, color=RGBColor(0xFF,0xFF,0xFF), bold=False)

# Quote
tb(s3,
   "\u201cThe best Nigerian talent isn\u2019t on Fiverr.\nIt\u2019s in someone\u2019s phone book.\nThe question is whose.\u201d",
   Inches(0.4), Inches(1.5), left_w - Inches(0.5), Inches(3.5),
   font_name="Georgia", font_size=15, italic=True, color=WHITE)

# Attribution
tb(s3, "— The problem we were built to solve",
   Inches(0.4), H - Inches(1), left_w - Inches(0.5), Inches(0.5),
   font_size=8, italic=True, color=RGBColor(0xFF,0xFF,0xFF))

# Right column diff items
diffs = [
    ("We brief ourselves before we brief anyone else",
     "One conversation about your business — not your creative spec. We understand the market, the competitive pressure, the internal politics. Then we translate it into a brief no freelancer would ever write."),
    ("Our roster is built on trust, not algorithms",
     "Every creative we place has worked with us before. We know how they behave under deadline pressure. We know what they're brilliant at and what they shouldn't touch. Platforms don't know that. We do."),
    ("We stay in until it ships",
     "We don't disappear after the match. We're in the thread, on the call, accountable for the outcome alongside the talent. Your brief is our brief until the final file lands."),
]

row_h = (H - Inches(0.5)) / len(diffs)
rx = left_w + Inches(0.3)
rw = W - left_w - Inches(0.5)

for i, (title, body) in enumerate(diffs):
    ry = Inches(0.25) + i * row_h

    # Divider line
    if i > 0:
        div = s3.shapes.add_shape(1, rx, ry, rw, Pt(0.5))
        div.fill.solid()
        div.fill.fore_color.rgb = RGBColor(0x1F,0x1F,0x1F)
        div.line.fill.background()

    # Arrow
    tb(s3, "→", rx, ry + Inches(0.15), Inches(0.3), Inches(0.4),
       font_size=13, color=RED)

    # Title
    tb(s3, title, rx + Inches(0.35), ry + Inches(0.15), rw - Inches(0.4), Inches(0.5),
       font_size=9, bold=True, color=WHITE)

    # Body
    tb(s3, body, rx + Inches(0.35), ry + Inches(0.65), rw - Inches(0.4), row_h - Inches(0.8),
       font_size=8, color=GREY6)

tb(s3, "03 / 07", W - Inches(1), H - Inches(0.3), Inches(0.8), Inches(0.25),
   font_size=7, color=RGBColor(0x55,0x55,0x55), align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4: WHO WE SERVE
# ─────────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
rect(s4, 0, 0, W, H, BONE)

half4 = W / 2

# Left copy
tb(s4, "WHO WE SERVE", Inches(0.4), Inches(0.55), half4 - Inches(0.5), Inches(0.4),
   font_size=7.5, bold=True, color=RED)

tb(s4, "Enterprise teams that can't afford to get creative wrong.",
   Inches(0.4), Inches(1.0), half4 - Inches(0.6), Inches(2.2),
   font_name="Georgia", font_size=22, color=INK)

tb(s4,
   "The Nigerian market is unforgiving. Consumers are sophisticated, competition is brutal, and attention is expensive. The brands that win are the ones where the business strategy and the creative execution are in the same conversation.",
   Inches(0.4), Inches(3.4), half4 - Inches(0.6), Inches(2.5),
   font_size=8.5, color=GREY6)

# Right segment grid
segments = [
    ("A", "CMOs & Marketing Directors",
     "Campaign concepting, brand refresh, content at speed — without losing control of the message or the money."),
    ("B", "Strategy & Innovation Teams",
     "Board-ready visual storytelling, pitch decks, and communication design that makes the numbers legible."),
    ("C", "HR & Internal Comms",
     "Employer branding and internal campaigns for teams that know their people are also watching the brand."),
    ("D", "Product & GTM Teams",
     "Launch assets, UX copy, and creative direction calibrated to a Nigerian consumer's expectations — not a template's."),
]

seg_h = H / 2
for i, (letter, title, body) in enumerate(segments):
    col = i % 2
    row = i // 2
    sx = half4 + col * (half4 / 2)
    sy = row * seg_h

    bg = BONE if (col + row) % 2 == 0 else RGBColor(0xEE,0xE7,0xD8)
    rect(s4, sx, sy, half4 / 2, seg_h, bg)

    # Ghost letter
    tb(s4, letter, sx + half4/2 - Inches(0.55), sy + Inches(0.05), Inches(0.55), Inches(0.9),
       font_name="Arial Black", font_size=44, bold=True, color=SAND, align=PP_ALIGN.RIGHT)

    tb(s4, title.upper(), sx + Inches(0.2), sy + Inches(0.2), half4/2 - Inches(0.35), Inches(0.55),
       font_size=7.5, bold=True, color=INK)

    tb(s4, body, sx + Inches(0.2), sy + Inches(0.85), half4/2 - Inches(0.35), seg_h - Inches(1.1),
       font_size=7.5, color=GREY6)

tb(s4, "04 / 07", W - Inches(1), H - Inches(0.3), Inches(0.8), Inches(0.25),
   font_size=7, color=GREY9, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5: TALENT
# ─────────────────────────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
rect(s5, 0, 0, W, H, DARK)

# Headline
tb(s5, "THE\nTALENT.", Inches(0.5), Inches(0.5), Inches(4.5), Inches(3),
   font_name="Arial Black", font_size=52, bold=True, color=WHITE)

# Red accent on "TALENT."
tb(s5, "TALENT.", Inches(0.5), Inches(1.8), Inches(4.5), Inches(1.2),
   font_name="Arial Black", font_size=52, bold=True, color=RED)

# Body copy
body_txt = (
    "A curated roster. Not a marketplace.\n\n"
    "Nigeria has extraordinary creative talent. The problem has never been scarcity — "
    "it's discoverability, accountability, and fit. We've spent years finding the people "
    "who consistently deliver under real conditions: tight timelines, complex stakeholders, "
    "and high-stakes briefs. You get their best work because they trust the environment "
    "we've built around them."
)
tb(s5, body_txt, Inches(7), Inches(0.5), Inches(5.8), Inches(4.5),
   font_size=9, color=GREY6)

# Tags
tags = [
    ("Brand Strategists", True), ("Creative Directors", False), ("Copywriters", False),
    ("Motion & Video Producers", True), ("UI/UX Designers", False), ("Visual Designers", False),
    ("Pitch Deck Specialists", False), ("Campaign Conceptors", True),
    ("Photographers", False), ("Content Strategists", False),
]

tx = Inches(0.5)
ty = H - Inches(1.8)
for tag_text, lit in tags:
    tag_w = Inches(len(tag_text) * 0.085 + 0.4)
    tag_h = Inches(0.35)
    fill = RED if lit else MID
    rect(s5, tx, ty, tag_w, tag_h, fill)
    tb(s5, tag_text.upper(), tx + Inches(0.12), ty + Inches(0.05),
       tag_w - Inches(0.15), tag_h - Inches(0.05),
       font_size=6.5, bold=True, color=WHITE if lit else DIM)
    tx += tag_w + Inches(0.12)
    if tx > W - Inches(2.5):
        tx = Inches(0.5)
        ty += tag_h + Inches(0.15)

tb(s5, "05 / 07", W - Inches(1), H - Inches(0.3), Inches(0.8), Inches(0.25),
   font_size=7, color=RGBColor(0x55,0x55,0x55), align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6: HOW IT WORKS
# ─────────────────────────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
rect(s6, 0, 0, W, H, BONE)

# Header bar
hbar6 = Inches(1.1)
rect(s6, 0, 0, W, hbar6, INK)
tb(s6, "How It Works", Inches(0.4), Inches(0.2), Inches(7), Inches(0.75),
   font_name="Arial Black", font_size=26, bold=True, color=WHITE)
tb(s6, "48-hour turnaround on every match", Inches(8.5), Inches(0.38), Inches(4), Inches(0.45),
   font_size=8.5, color=RED, align=PP_ALIGN.RIGHT)

steps = [
    ("1", "Brief Us",
     "One call or a short form. We want to understand the business challenge — not just the deliverable. What's at stake? What's failed before?",
     "Same day"),
    ("2", "We Match",
     "2–3 curated options. Pre-briefed, pre-vetted, with a clear recommended approach. No shortlist of strangers — people we'd stake our name on.",
     "Within 48 hours"),
    ("3", "You Approve",
     "Pick your match. We handle contracts, onboarding, and kick-off. You brief once. Work starts immediately.",
     "Day 3"),
    ("4", "We Deliver",
     "Milestones, feedback loops, and final sign-off. We stay in the thread from first brief to final file. The accountability doesn't leave when the talent arrives.",
     "On deadline"),
]

step_w = W / 4
for i, (num, title, body, time_label) in enumerate(steps):
    sx = i * step_w
    sy = hbar6
    sh = H - hbar6

    bg = BONE if i % 2 == 0 else RGBColor(0xEE,0xE7,0xD8)
    rect(s6, sx, sy, step_w, sh, bg)

    # Step number
    tb(s6, num, sx + Inches(0.2), sy + Inches(0.3), step_w - Inches(0.3), Inches(0.85),
       font_name="Arial Black", font_size=38, bold=True, color=RED)

    # Title
    tb(s6, title.upper(), sx + Inches(0.2), sy + Inches(1.2), step_w - Inches(0.3), Inches(0.45),
       font_size=8.5, bold=True, color=INK)

    # Body
    tb(s6, body, sx + Inches(0.2), sy + Inches(1.7), step_w - Inches(0.35), sh - Inches(2.9),
       font_size=8, color=GREY6)

    # Time tag
    tb(s6, time_label, sx + Inches(0.2), H - Inches(0.65), step_w - Inches(0.3), Inches(0.4),
       font_size=7.5, bold=True, color=RED)

    # Divider
    if i > 0:
        div = s6.shapes.add_shape(1, sx, sy, Pt(0.5), sh)
        div.fill.solid()
        div.fill.fore_color.rgb = SAND
        div.line.fill.background()

tb(s6, "06 / 07", W - Inches(1), H - Inches(0.3), Inches(0.8), Inches(0.25),
   font_size=7, color=GREY9, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7: CLOSE
# ─────────────────────────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
rect(s7, 0, 0, W, H, INK)

left7 = W * 58 // 100
right7 = W - left7

# Left panel
tb(s7, "START HERE", Inches(0.5), Inches(0.5), left7 - Inches(0.7), Inches(0.4),
   font_size=7.5, bold=True, color=RED)

tb(s7,
   "Your next campaign deserves a partner who understands this market.",
   Inches(0.5), Inches(1.1), left7 - Inches(0.7), Inches(2.4),
   font_name="Georgia", font_size=26, color=WHITE)

tb(s7,
   "We're not a platform. We're not an agency. We're the people who know the right person "
   "for your brief — and who make sure they deliver. Lagos, Abuja, Port Harcourt. One brief away.",
   Inches(0.5), Inches(4.0), left7 - Inches(0.7), Inches(2.5),
   font_size=9, color=GREY6)

# Right red panel
rect(s7, left7, 0, right7, H, RED)

# CTA label
tb(s7, "GET IN TOUCH", left7 + Inches(0.3), Inches(0.55), right7 - Inches(0.4), Inches(0.4),
   font_size=7, color=RGBColor(0xFF,0xFF,0xFF))

# CTA main
tb(s7, "Brief Us\nToday.", left7 + Inches(0.3), Inches(1.3), right7 - Inches(0.4), Inches(2.5),
   font_name="Arial Black", font_size=36, bold=True, color=WHITE)

# Contact details
for j, line in enumerate(["hello@commcorp.co", "www.commcorp.co", "Lagos · Abuja · Port Harcourt"]):
    ly = Inches(5.2) + j * Inches(0.5)
    col = RGBColor(0xFF,0xFF,0xFF) if j < 2 else RGBColor(0xFF,0xAA,0xAA)
    tb(s7, line, left7 + Inches(0.3), ly, right7 - Inches(0.4), Inches(0.45),
       font_size=9 if j < 2 else 7.5, color=col)

tb(s7, "07 / 07", W - Inches(0.8), H - Inches(0.3), Inches(0.7), Inches(0.25),
   font_size=7, color=RGBColor(0x55,0x55,0x55), align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "commcorp_deck.pptx")
prs.save(out)
print(f"✅ Saved: {out}")
